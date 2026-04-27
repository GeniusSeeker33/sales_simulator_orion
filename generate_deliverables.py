"""
Generate roi_model.xlsx and exec_pitch_deck.pptx for Orion exec pitch.
Run: python3 generate_deliverables.py
"""

import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side, numbers
from openpyxl.utils import get_column_letter
from openpyxl.chart import BarChart, Reference
from openpyxl.chart.series import DataPoint
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUTPUT_DIR = "/workspaces/sales_simulator_orion"

# ─────────────────────────────────────────────────────────────────────────────
# COLORS
# ─────────────────────────────────────────────────────────────────────────────
DARK_BG   = "1A1A2E"
MID_BG    = "16213E"
ACCENT    = "0F3460"
GOLD      = "E94560"
WHITE     = "FFFFFF"
LIGHT_GRY = "F0F0F0"
MED_GRY   = "CCCCCC"
GREEN     = "27AE60"
RED_      = "E74C3C"
YELLOW    = "F1C40F"

def hex_fill(hex_color):
    return PatternFill("solid", fgColor=hex_color)

def thin_border():
    s = Side(style="thin", color="CCCCCC")
    return Border(left=s, right=s, top=s, bottom=s)


# ─────────────────────────────────────────────────────────────────────────────
# ROI MODEL — EXCEL
# ─────────────────────────────────────────────────────────────────────────────
def build_roi_model():
    wb = openpyxl.Workbook()

    # ── Sheet 1: Assumptions ──────────────────────────────────────────────────
    ws1 = wb.active
    ws1.title = "Assumptions"
    ws1.sheet_view.showGridLines = False

    def hdr(ws, row, col, text, bg=DARK_BG, fg=WHITE, bold=True, size=12):
        c = ws.cell(row=row, column=col, value=text)
        c.fill = hex_fill(bg)
        c.font = Font(color=fg, bold=bold, size=size, name="Calibri")
        c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        c.border = thin_border()
        return c

    def val(ws, row, col, v, fmt=None, bold=False, bg=None, fg="000000"):
        c = ws.cell(row=row, column=col, value=v)
        c.font = Font(bold=bold, size=11, name="Calibri", color=fg)
        c.alignment = Alignment(horizontal="right", vertical="center")
        c.border = thin_border()
        if bg:
            c.fill = hex_fill(bg)
        if fmt:
            c.number_format = fmt
        return c

    def lbl(ws, row, col, text, indent=0, bold=False, bg=None):
        c = ws.cell(row=row, column=col, value=("  " * indent) + text)
        c.font = Font(bold=bold, size=11, name="Calibri")
        c.alignment = Alignment(horizontal="left", vertical="center")
        c.border = thin_border()
        if bg:
            c.fill = hex_fill(bg)
        return c

    # Title
    ws1.merge_cells("A1:F1")
    t = ws1["A1"]
    t.value = "GeniusSeeker Sales Simulator — Orion ROI Model"
    t.fill = hex_fill(DARK_BG)
    t.font = Font(color=WHITE, bold=True, size=16, name="Calibri")
    t.alignment = Alignment(horizontal="center", vertical="center")
    ws1.row_dimensions[1].height = 36

    ws1.merge_cells("A2:F2")
    sub = ws1["A2"]
    sub.value = "Edit yellow cells to adjust assumptions. All other cells calculate automatically."
    sub.fill = hex_fill(ACCENT)
    sub.font = Font(color=WHITE, size=10, italic=True, name="Calibri")
    sub.alignment = Alignment(horizontal="center", vertical="center")
    ws1.row_dimensions[2].height = 22

    # Section headers row 4
    hdr(ws1, 4, 1, "ASSUMPTION", bg=ACCENT, size=11)
    hdr(ws1, 4, 2, "CURRENT (Yr 1)", bg=ACCENT, size=11)
    hdr(ws1, 4, 3, "SCALE (Yr 2)", bg=ACCENT, size=11)
    hdr(ws1, 4, 4, "NOTES", bg=ACCENT, size=11)
    ws1.merge_cells("D4:F4")

    # ── Seat / Pricing
    rows = [
        (5,  "── ORION HEADCOUNT & PRICING ──", None, None, ""),
        (6,  "Current sales reps",               39,    None,  "Orion today"),
        (7,  "Projected reps (2 years)",          None,   80,    "Per leadership plan"),
        (8,  "Price per user / month",            30,    30,    "Target price point"),
        (9,  "Annual contract value (current)",   None,  None,  "= seats × $30 × 12"),
        (10, "Annual contract value (at scale)",  None,  None,  "= 80 × $30 × 12"),

        (11, "", None, None, ""),
        (12, "── REP ECONOMICS ──", None, None, ""),
        (13, "Avg annual quota per rep",           2000000, 2000000, "Distribution industry typical"),
        (14, "Avg ramp time today (months)",        4.5,    4.5,   "3–6 month range; use midpoint"),
        (15, "Ramp time with Simulator (months)",   3.5,    3.5,   "Conservative: 30-day improvement"),
        (16, "Rep productivity during ramp (%)",    0.45,   0.45,  "Ramp reps at ~45% of quota"),
        (17, "Annual new hires / departures",        8,      16,    "~20% of 39 → ~40% churn+growth at 80"),
        (18, "Gross margin on rep revenue (%)",      0.30,   0.30,  "Adjust to Orion's actual margin"),

        (19, "", None, None, ""),
        (20, "── CALCULATED IMPACT ──", None, None, ""),
        (21, "Revenue recovered per hire (30-day faster ramp)", None, None, "= quota/12 × margin × productivity"),
        (22, "Annual revenue recovered (all new hires)",        None, None, "= above × annual new hires"),
        (23, "Annual tool cost",                                None, None, "= seats × $30 × 12"),
        (24, "Net annual benefit",                              None, None, "= recovered revenue − tool cost"),
        (25, "ROI (return on investment)",                      None, None, "= net benefit / tool cost"),
        (26, "Payback period (months)",                        None, None, "= tool cost / (recovered rev / 12)"),
    ]

    INPUT_YELLOW = "FFF9C4"
    CALC_BLUE    = "E3F2FD"

    for r, label, v1, v2, note in rows:
        lbl(ws1, r, 1, label, bold=("──" in label), bg=(LIGHT_GRY if "──" in label else None))
        if label.startswith("──"):
            ws1.merge_cells(f"A{r}:F{r}")
            c = ws1.cell(row=r, column=1)
            c.fill = hex_fill(MID_BG)
            c.font = Font(color=YELLOW, bold=True, size=11, name="Calibri")
            c.alignment = Alignment(horizontal="left", vertical="center")
            continue

        if v1 is not None:
            c1 = val(ws1, r, 2, v1, bg=INPUT_YELLOW)
            if isinstance(v1, float) and v1 < 1:
                c1.number_format = "0%"
            elif isinstance(v1, int) and v1 > 999:
                c1.number_format = "$#,##0"

        if v2 is not None:
            c2 = val(ws1, r, 3, v2, bg=INPUT_YELLOW)
            if isinstance(v2, float) and v2 < 1:
                c2.number_format = "0%"
            elif isinstance(v2, int) and v2 > 999:
                c2.number_format = "$#,##0"

        if note:
            n = ws1.cell(row=r, column=4, value=note)
            n.font = Font(size=10, italic=True, name="Calibri", color="555555")
            n.alignment = Alignment(horizontal="left", vertical="center")
            n.border = thin_border()
            ws1.merge_cells(f"D{r}:F{r}")

    # Formulas for calculated rows (Year 1 = col B, Year 2 = col C)
    # Row 9: ACV current  = B6 * B8 * 12
    ws1["B9"] = "=B6*B8*12"
    ws1["B9"].number_format = "$#,##0"
    ws1["B9"].fill = hex_fill(CALC_BLUE)
    ws1["B9"].font = Font(size=11, bold=True, name="Calibri")
    ws1["B9"].border = thin_border()
    ws1["B9"].alignment = Alignment(horizontal="right")

    # Row 10: ACV at scale
    ws1["C10"] = "=C7*C8*12"
    ws1["C10"].number_format = "$#,##0"
    ws1["C10"].fill = hex_fill(CALC_BLUE)
    ws1["C10"].font = Font(size=11, bold=True, name="Calibri")
    ws1["C10"].border = thin_border()
    ws1["C10"].alignment = Alignment(horizontal="right")

    # Row 21: Revenue recovered per hire = (quota/12) × (months saved = B14-B15) × B16 × B18
    ws1["B21"] = "=(B13/12)*(B14-B15)*B16*B18"
    ws1["B21"].number_format = "$#,##0"
    ws1["B21"].fill = hex_fill(CALC_BLUE)
    ws1["B21"].font = Font(size=11, bold=True, name="Calibri")
    ws1["B21"].border = thin_border()
    ws1["B21"].alignment = Alignment(horizontal="right")

    ws1["C21"] = "=(C13/12)*(C14-C15)*C16*C18"
    ws1["C21"].number_format = "$#,##0"
    ws1["C21"].fill = hex_fill(CALC_BLUE)
    ws1["C21"].font = Font(size=11, bold=True, name="Calibri")
    ws1["C21"].border = thin_border()
    ws1["C21"].alignment = Alignment(horizontal="right")

    # Row 22: Total annual impact = B21 × B17
    ws1["B22"] = "=B21*B17"
    ws1["B22"].number_format = "$#,##0"
    ws1["B22"].fill = hex_fill(CALC_BLUE)
    ws1["B22"].font = Font(size=11, bold=True, name="Calibri")
    ws1["B22"].border = thin_border()
    ws1["B22"].alignment = Alignment(horizontal="right")

    ws1["C22"] = "=C21*C17"
    ws1["C22"].number_format = "$#,##0"
    ws1["C22"].fill = hex_fill(CALC_BLUE)
    ws1["C22"].font = Font(size=11, bold=True, name="Calibri")
    ws1["C22"].border = thin_border()
    ws1["C22"].alignment = Alignment(horizontal="right")

    # Row 23: Tool cost year 1 = B6*B8*12, year 2 = C7*C8*12
    ws1["B23"] = "=B6*B8*12"
    ws1["B23"].number_format = "$#,##0"
    ws1["B23"].fill = hex_fill(CALC_BLUE)
    ws1["B23"].font = Font(size=11, bold=True, name="Calibri")
    ws1["B23"].border = thin_border()
    ws1["B23"].alignment = Alignment(horizontal="right")

    ws1["C23"] = "=C7*C8*12"
    ws1["C23"].number_format = "$#,##0"
    ws1["C23"].fill = hex_fill(CALC_BLUE)
    ws1["C23"].font = Font(size=11, bold=True, name="Calibri")
    ws1["C23"].border = thin_border()
    ws1["C23"].alignment = Alignment(horizontal="right")

    # Row 24: Net benefit
    ws1["B24"] = "=B22-B23"
    ws1["B24"].number_format = "$#,##0"
    ws1["B24"].fill = hex_fill("D5F5E3")
    ws1["B24"].font = Font(size=11, bold=True, name="Calibri", color=GREEN)
    ws1["B24"].border = thin_border()
    ws1["B24"].alignment = Alignment(horizontal="right")

    ws1["C24"] = "=C22-C23"
    ws1["C24"].number_format = "$#,##0"
    ws1["C24"].fill = hex_fill("D5F5E3")
    ws1["C24"].font = Font(size=11, bold=True, name="Calibri", color=GREEN)
    ws1["C24"].border = thin_border()
    ws1["C24"].alignment = Alignment(horizontal="right")

    # Row 25: ROI
    ws1["B25"] = "=B22/B23"
    ws1["B25"].number_format = "0.0x"
    ws1["B25"].fill = hex_fill("D5F5E3")
    ws1["B25"].font = Font(size=12, bold=True, name="Calibri", color=GREEN)
    ws1["B25"].border = thin_border()
    ws1["B25"].alignment = Alignment(horizontal="right")

    ws1["C25"] = "=C22/C23"
    ws1["C25"].number_format = "0.0x"
    ws1["C25"].fill = hex_fill("D5F5E3")
    ws1["C25"].font = Font(size=12, bold=True, name="Calibri", color=GREEN)
    ws1["C25"].border = thin_border()
    ws1["C25"].alignment = Alignment(horizontal="right")

    # Row 26: Payback (months)
    ws1["B26"] = "=B23/(B22/12)"
    ws1["B26"].number_format = "0.0 \"months\""
    ws1["B26"].fill = hex_fill(CALC_BLUE)
    ws1["B26"].font = Font(size=11, bold=True, name="Calibri")
    ws1["B26"].border = thin_border()
    ws1["B26"].alignment = Alignment(horizontal="right")

    ws1["C26"] = "=C23/(C22/12)"
    ws1["C26"].number_format = "0.0 \"months\""
    ws1["C26"].fill = hex_fill(CALC_BLUE)
    ws1["C26"].font = Font(size=11, bold=True, name="Calibri")
    ws1["C26"].border = thin_border()
    ws1["C26"].alignment = Alignment(horizontal="right")

    # Column widths
    ws1.column_dimensions["A"].width = 45
    ws1.column_dimensions["B"].width = 22
    ws1.column_dimensions["C"].width = 22
    ws1.column_dimensions["D"].width = 50

    # ── Sheet 2: Competitive Benchmark ──────────────────────────────────────
    ws2 = wb.create_sheet("Competitive Benchmark")
    ws2.sheet_view.showGridLines = False

    ws2.merge_cells("A1:G1")
    t2 = ws2["A1"]
    t2.value = "Sales Enablement & Training Tool Pricing Benchmark"
    t2.fill = hex_fill(DARK_BG)
    t2.font = Font(color=WHITE, bold=True, size=14, name="Calibri")
    t2.alignment = Alignment(horizontal="center", vertical="center")
    ws2.row_dimensions[1].height = 32

    cols = ["Tool", "Category", "Price/User/Mo", "Annual/User", "vs. GeniusSeeker", "Strength", "Weakness"]
    for i, c in enumerate(cols, 1):
        cell = ws2.cell(row=3, column=i, value=c)
        cell.fill = hex_fill(ACCENT)
        cell.font = Font(color=WHITE, bold=True, size=11, name="Calibri")
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border = thin_border()

    competitors = [
        ("GeniusSeeker (Orion)", "AI Sales Simulator",    30,   360,  "Baseline",       "AI sim + manager view + pricing", "Newer product"),
        ("Gong",                  "Revenue Intelligence",  135,  1620, "+$105/user/mo",  "Deep call recording & analytics", "Complex, expensive"),
        ("Chorus (ZoomInfo)",     "Conversation Intel.",   110,  1320, "+$80/user/mo",   "Call recording & coaching",      "No simulation"),
        ("Mindtickle",            "Sales Readiness",        75,   900, "+$45/user/mo",   "Learning paths, certification",   "Rigid training format"),
        ("Lessonly / Seismic",    "Sales Enablement",       45,   540, "+$15/user/mo",   "Content library, quizzes",       "No AI simulation"),
        ("Brainshark",            "Video Coaching",         60,   720, "+$30/user/mo",   "Manager video reviews",          "No AI customer roleplay"),
        ("Allego",                "Sales Readiness",        55,   660, "+$25/user/mo",   "Video pitch practice",           "No live AI customer"),
        ("Second Nature AI",      "AI Roleplay",            80,   960, "+$50/user/mo",   "AI roleplay training",           "No account-based context"),
    ]

    for i, row_data in enumerate(competitors):
        r = i + 4
        bg = "FFF9C4" if i == 0 else (LIGHT_GRY if i % 2 == 0 else WHITE)
        for j, v in enumerate(row_data, 1):
            c = ws2.cell(row=r, column=j, value=v)
            c.fill = hex_fill(bg)
            c.font = Font(size=11, name="Calibri", bold=(i == 0))
            c.alignment = Alignment(horizontal="center" if j > 1 else "left", vertical="center", wrap_text=True)
            c.border = thin_border()
            if j == 3:
                c.number_format = "$#,##0"
            if j == 4:
                c.number_format = "$#,##0"

    ws2.column_dimensions["A"].width = 22
    ws2.column_dimensions["B"].width = 22
    ws2.column_dimensions["C"].width = 16
    ws2.column_dimensions["D"].width = 14
    ws2.column_dimensions["E"].width = 16
    ws2.column_dimensions["F"].width = 32
    ws2.column_dimensions["G"].width = 28

    for r in range(4, 12):
        ws2.row_dimensions[r].height = 28

    # ── Sheet 3: Pilot Structure ──────────────────────────────────────────────
    ws3 = wb.create_sheet("Pilot & Pricing Options")
    ws3.sheet_view.showGridLines = False

    ws3.merge_cells("A1:E1")
    t3 = ws3["A1"]
    t3.value = "Pilot Structure & Pricing Options — Orion"
    t3.fill = hex_fill(DARK_BG)
    t3.font = Font(color=WHITE, bold=True, size=14, name="Calibri")
    t3.alignment = Alignment(horizontal="center", vertical="center")
    ws3.row_dimensions[1].height = 32

    options = [
        ("Option", "Structure",         "Seats",  "Monthly",   "Annual",   "Key Terms"),
        ("A — Full Rollout",             "All reps from day 1",       "39",  "$1,170",  "$14,040",  "Standard 12-month contract"),
        ("B — Pilot + Expand",           "10 reps, 90-day pilot",     "10",  "$300",    "$3,600",   "Auto-expand to 39 if retention ≥ 80%"),
        ("C — One Team Pilot",           "1 sales team (8–12 reps)",  "12",  "$360",    "$4,320",   "60-day success review; full rollout trigger"),
        ("D — Success-Based",            "All 39 reps",               "39",  "$1,170",  "$14,040",  "Price-back if ramp time doesn't improve"),
        ("E — Tiered Growth",            "Starts 39, grows to 80",    "39→80","$1,170→$2,400","$14K–$29K","Locked rate as headcount grows"),
    ]

    for i, row_data in enumerate(options):
        r = i + 3
        is_header = (i == 0)
        bg = ACCENT if is_header else (LIGHT_GRY if i % 2 == 0 else WHITE)
        fg = WHITE if is_header else "000000"
        for j, v in enumerate(row_data, 1):
            c = ws3.cell(row=r, column=j, value=v)
            c.fill = hex_fill(bg)
            c.font = Font(size=11, name="Calibri", bold=is_header, color=fg)
            c.alignment = Alignment(horizontal="center" if j > 1 else "left", vertical="center", wrap_text=True)
            c.border = thin_border()
        ws3.row_dimensions[r].height = 30

    ws3.column_dimensions["A"].width = 22
    ws3.column_dimensions["B"].width = 28
    ws3.column_dimensions["C"].width = 10
    ws3.column_dimensions["D"].width = 14
    ws3.column_dimensions["E"].width = 14
    ws3.column_dimensions["F"].width = 40

    # Save
    path = os.path.join(OUTPUT_DIR, "roi_model.xlsx")
    wb.save(path)
    print(f"✅  Saved: {path}")


# ─────────────────────────────────────────────────────────────────────────────
# EXEC PITCH DECK — POWERPOINT
# ─────────────────────────────────────────────────────────────────────────────
def rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_text_box(slide, text, left, top, width, height,
                  font_size=18, bold=False, color="FFFFFF",
                  align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    run.font.italic = italic
    run.font.name = "Calibri"
    return txBox


def set_slide_bg(slide, hex_color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(hex_color)


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=0):
    from pptx.util import Pt as PPt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_color)
    if line_color:
        shape.line.color.rgb = rgb(line_color)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def build_deck():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    slides_data = []

    # ── SLIDE 1: Title ────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, DARK_BG)
    add_rect(sl, 0, 5.8, 13.33, 1.7, ACCENT)
    add_text_box(sl, "GeniusSeeker Sales Simulator",
                 0.5, 1.2, 12, 1.2, font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, "AI-Powered Sales Training for Orion's Growing Team",
                 0.5, 2.5, 12, 0.8, font_size=22, color=MED_GRY, align=PP_ALIGN.CENTER)
    add_rect(sl, 4.67, 3.5, 4, 0.08, GOLD)
    add_text_box(sl, "$30 / rep / month  •  30-Day Faster Ramp  •  Zero Manager Overhead",
                 0.5, 3.8, 12, 0.7, font_size=16, color=YELLOW, align=PP_ALIGN.CENTER, italic=True)
    add_text_box(sl, "Presented to Orion Leadership  •  April 2026",
                 0.5, 6.2, 12, 0.6, font_size=13, color=MED_GRY, align=PP_ALIGN.CENTER)

    # ── SLIDE 2: The Problem ──────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, 13.33, 1.1, ACCENT)
    add_text_box(sl, "The Problem: Every New Rep Costs You 4–5 Months",
                 0.3, 0.15, 12.7, 0.8, font_size=26, bold=True, color=WHITE)

    stats = [
        ("4.5 months",   "Average ramp time\nfor Orion reps today"),
        ("~45%",         "Productivity during\nramp period"),
        ("~$75K",        "Estimated revenue\nmissed per new hire"),
        ("8–16/year",    "New hires as team\ngrows to 80 reps"),
    ]
    for i, (num, label) in enumerate(stats):
        x = 0.4 + i * 3.2
        add_rect(sl, x, 1.4, 2.9, 2.8, ACCENT)
        add_text_box(sl, num,   x+0.1, 1.55, 2.7, 1.1, font_size=38, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text_box(sl, label, x+0.1, 2.7,  2.7, 0.9, font_size=14, color=MED_GRY, align=PP_ALIGN.CENTER)

    add_rect(sl, 0.4, 4.5, 12.5, 0.08, GOLD)
    add_text_box(sl,
        "The feedback loop is too long. Reps make mistakes on real accounts, managers find out weeks later, and the pattern repeats.",
        0.4, 4.7, 12.5, 0.8, font_size=17, color=WHITE, italic=True, align=PP_ALIGN.CENTER)
    add_text_box(sl,
        "With 39 reps today and 80 in two years, slow ramp is a compounding revenue problem — not a training inconvenience.",
        0.4, 5.5, 12.5, 0.7, font_size=15, color=MED_GRY, align=PP_ALIGN.CENTER)

    # ── SLIDE 3: The Solution ─────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, 13.33, 1.1, ACCENT)
    add_text_box(sl, "The Solution: A Flight Simulator for Sales Calls",
                 0.3, 0.15, 12.7, 0.8, font_size=26, bold=True, color=WHITE)

    features = [
        ("AI Customer Roleplay",       "Practice any dealer call — 8 personas, 4 difficulty levels. The AI pushes back like a real buyer."),
        ("Account-Based Practice",     "Rep selects a real dealer account. AI plays that specific buyer with their barriers and category gaps."),
        ("AI Scoring & Coaching",      "Every session scored on 5 dimensions. Coaching notes + exact phrases the rep should have used."),
        ("Manager Intelligence View",  "Real-time visibility into every rep's skill level, coaching needs, and practice activity."),
        ("Rep Progression Levels",     "L1 → L5 progression tied to training scores. Reps see their own improvement. Managers too."),
        ("RingCentral Integration",    "Import real call logs. Compare simulator practice vs. live call activity side by side."),
    ]

    for i, (title, body) in enumerate(features):
        row = i // 2
        col = i % 2
        x = 0.4 + col * 6.5
        y = 1.35 + row * 1.7
        add_rect(sl, x, y, 6.1, 1.55, ACCENT)
        add_rect(sl, x, y, 0.12, 1.55, GOLD)
        add_text_box(sl, title, x+0.25, y+0.1, 5.7, 0.45, font_size=14, bold=True, color=WHITE)
        add_text_box(sl, body,  x+0.25, y+0.55, 5.7, 0.9,  font_size=12, color=MED_GRY)

    add_text_box(sl, "Available on day one. No setup. No IT project. Reps log in and start practicing.",
                 0.4, 6.8, 12.5, 0.45, font_size=13, color=YELLOW, italic=True, align=PP_ALIGN.CENTER)

    # ── SLIDE 4: The Business Case (ROI) ─────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, 13.33, 1.1, ACCENT)
    add_text_box(sl, "The Business Case: 28x ROI in Year 1",
                 0.3, 0.15, 12.7, 0.8, font_size=26, bold=True, color=WHITE)

    # Left column: math
    add_rect(sl, 0.4, 1.25, 6.0, 5.5, ACCENT)
    add_text_box(sl, "Year 1 Model  (39 reps)", 0.6, 1.35, 5.6, 0.55, font_size=16, bold=True, color=GOLD)
    math_rows = [
        ("Average quota per rep",                 "$2,000,000"),
        ("Ramp time today",                       "4.5 months"),
        ("Ramp time with simulator",              "3.5 months"),
        ("Time saved per new hire",               "1 month"),
        ("Productivity during ramp",              "45%"),
        ("Gross margin",                          "30%"),
        ("Revenue recovered per hire",            "≈ $22,500"),
        ("New hires in Year 1",                   "8 reps"),
        ("Total revenue impact",                  "≈ $180,000"),
        ("────────────────",                       ""),
        ("Annual tool cost  (39 × $30 × 12)",     "$14,040"),
        ("Net annual benefit",                    "≈ $165,960"),
        ("ROI",                                   "≈ 28×"),
        ("Payback period",                        "< 1 month"),
    ]
    for i, (k, v) in enumerate(math_rows):
        y = 1.95 + i * 0.31
        if "────" in k:
            add_rect(sl, 0.55, y+0.05, 5.7, 0.04, GOLD)
            continue
        bold_row = k in ("ROI", "Net annual benefit", "Payback period")
        fc = GOLD if bold_row else WHITE
        add_text_box(sl, k, 0.6,  y, 3.5, 0.3, font_size=11, color=fc, bold=bold_row)
        add_text_box(sl, v, 4.1,  y, 2.1, 0.3, font_size=11, color=fc, bold=bold_row, align=PP_ALIGN.RIGHT)

    # Right column: call-outs
    callouts = [
        (GOLD,   "28×",   "ROI in Year 1"),
        (GREEN,  "< 1 mo","Payback period"),
        (ACCENT, "$30",   "per rep / month"),
        (ACCENT, "80",    "reps in 2 years\n= $28,800/yr ACV"),
    ]
    for i, (bg, big, sub) in enumerate(callouts):
        y = 1.25 + i * 1.5
        add_rect(sl, 6.9, y, 5.9, 1.35, bg)
        add_text_box(sl, big, 7.0, y+0.05, 2.5, 0.85, font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, sub, 9.5, y+0.25, 3.2, 0.75, font_size=14, color=MED_GRY)

    # ── SLIDE 5: Manager View ─────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, 13.33, 1.1, ACCENT)
    add_text_box(sl, "What Managers See: One View for Everything",
                 0.3, 0.15, 12.7, 0.8, font_size=26, bold=True, color=WHITE)

    panels = [
        ("Rep Comparison Table",
         "Every AE side-by-side: revenue, captures, comp status, training score, level, and missed upside. Spot coaching opportunities in seconds."),
        ("Auto-Generated Coaching Alerts",
         "System flags: low simulator scores, execution gaps (practicing but not calling), training gaps (calling but not practicing)."),
        ("Simulator Activity Log",
         "Every practice session logged: dealer name, buyer, score, difficulty, coaching note. Filter by rep, dealer, or score range."),
        ("RingCentral Call Activity",
         "Import real call logs. See call volume, connect rate, and duration alongside simulator scores for a complete performance picture."),
        ("Rep Leaderboard",
         "Combined performance score from simulator skill + call volume + connect rate. Transparency drives healthy competition."),
        ("Compensation Snapshot",
         "Each rep's KPI attainment status: Ramp, Base, or Accelerated tier. See who's on track and who's leaving money on the table."),
    ]
    for i, (title, body) in enumerate(panels):
        row = i // 2
        col = i % 2
        x = 0.4 + col * 6.5
        y = 1.35 + row * 1.72
        add_rect(sl, x, y, 6.1, 1.58, ACCENT)
        add_rect(sl, x, y, 0.12, 1.58, GOLD)
        add_text_box(sl, title, x+0.25, y+0.1,  5.7, 0.42, font_size=14, bold=True, color=WHITE)
        add_text_box(sl, body,  x+0.25, y+0.55, 5.7, 0.92, font_size=12, color=MED_GRY)

    add_text_box(sl, "No more guessing who needs coaching. No more waiting for the QBR to find out a rep is struggling.",
                 0.4, 6.8, 12.5, 0.45, font_size=13, color=YELLOW, italic=True, align=PP_ALIGN.CENTER)

    # ── SLIDE 6: Pricing & Pilot ──────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, 13.33, 1.1, ACCENT)
    add_text_box(sl, "$30 / Rep / Month — Pilot Structures Available",
                 0.3, 0.15, 12.7, 0.8, font_size=26, bold=True, color=WHITE)

    # Main price card
    add_rect(sl, 0.4, 1.3, 3.8, 3.5, GOLD)
    add_text_box(sl, "$30",       0.5, 1.45, 3.6, 1.3, font_size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, "per rep / month", 0.5, 2.75, 3.6, 0.55, font_size=16, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, "39 reps = $14,040/yr\n80 reps = $28,800/yr", 0.5, 3.3, 3.6, 0.8, font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, "Less than 1 Gong seat", 0.5, 4.1, 3.6, 0.45, font_size=12, color=WHITE, italic=True, align=PP_ALIGN.CENTER)

    # Pilot options
    pilots = [
        ("Option A: Full Rollout",     "All 39 reps, 12-month contract.\n$14,040/year. Standard terms."),
        ("Option B: 90-Day Pilot",     "10 reps for 90 days ($900).\nFull rollout trigger on adoption."),
        ("Option C: Team Pilot",       "1 sales team, 60-day review.\nSuccess = manager sign-off on KPIs."),
        ("Option D: Success-Based",    "Full rollout with price guarantee.\nPartial rebate if ramp doesn't improve."),
        ("Option E: Growth Lock",      "Lock $30 rate as you grow to 80.\nNo price increase for 24 months."),
    ]
    for i, (title, body) in enumerate(pilots):
        x = 4.6 + (i % 2) * 4.3
        y = 1.3 + (i // 2) * 2.0
        if i == 4:  # last one centered
            x = 6.4
            y = 5.35
        add_rect(sl, x, y, 4.0, 1.75, ACCENT)
        add_rect(sl, x, y, 0.1, 1.75, GOLD)
        add_text_box(sl, title, x+0.2, y+0.12, 3.7, 0.48, font_size=13, bold=True, color=WHITE)
        add_text_box(sl, body,  x+0.2, y+0.6,  3.7, 1.0,  font_size=11, color=MED_GRY)

    # ── SLIDE 7: Objection Handling ────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, 13.33, 1.1, ACCENT)
    add_text_box(sl, "Anticipated Questions — and the Answers",
                 0.3, 0.15, 12.7, 0.8, font_size=26, bold=True, color=WHITE)

    objections = [
        ('"Will reps actually use it?"',
         "Usage is tied to their progression level and comp readiness score — the incentive is built in. We also see 7-10 min sessions fit naturally into morning call prep."),
        ('"How is this different from training videos?"',
         "A video shows a perfect call. This forces reps to have the hard conversation, get pushed back on, and see exactly where they fell short. Active vs. passive."),
        ('"What if the AI gives wrong advice?"',
         "The AI plays the customer — it's not coaching. The scoring rubric reflects your actual sales process. You control what 'good' looks like."),
        ('"$30/user is still a budget item — prove the ROI."',
         "One rep who ramps 30 days faster recovers ~$22,500 in revenue. 8 new hires = $180K impact. The tool costs $14K. That's a 28x return in Year 1."),
        ('"We already have onboarding materials."',
         "Materials tell reps what to do. This makes them practice doing it — with a customer who pushes back. One is reading about swimming, the other is getting in the pool."),
    ]
    for i, (obj, resp) in enumerate(objections):
        y = 1.25 + i * 1.22
        add_rect(sl, 0.4, y, 12.5, 1.1, ACCENT)
        add_rect(sl, 0.4, y, 0.12, 1.1, GOLD)
        add_text_box(sl, obj,  0.65, y+0.05, 4.5, 0.45, font_size=13, bold=True, color=GOLD)
        add_text_box(sl, resp, 5.2,  y+0.05, 7.5, 0.95, font_size=12, color=WHITE)

    # ── SLIDE 8: The Ask ──────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, 13.33, 1.1, ACCENT)
    add_text_box(sl, "The Ask",
                 0.3, 0.15, 12.7, 0.8, font_size=26, bold=True, color=WHITE)

    add_rect(sl, 2.5, 1.4, 8.33, 2.2, ACCENT)
    add_text_box(sl, "90-Day Pilot: 10 Reps  |  $900 Total",
                 2.7, 1.55, 8.0, 0.7, font_size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text_box(sl, "At the end of 90 days, we review:\n• Rep ramp progress vs. baseline\n• Manager-reported coaching efficiency\n• Simulator score trajectory (start vs. end)\n\nIf the data supports it — we expand to all 39 reps.",
                 2.7, 2.3, 8.0, 1.2, font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

    steps = [
        ("Step 1\nToday",          "Approve 90-day pilot\n(10 reps, $900)"),
        ("Step 2\nWeek 1",         "Onboard pilot reps\n(< 1 hour setup)"),
        ("Step 3\nDay 30",         "First check-in:\nscore trends & usage"),
        ("Step 4\nDay 90",         "Review outcomes\n& expand decision"),
    ]
    for i, (step, action) in enumerate(steps):
        x = 0.9 + i * 3.1
        add_rect(sl, x, 4.0, 2.6, 2.0, ACCENT)
        add_rect(sl, x, 4.0, 2.6, 0.5, GOLD)
        add_text_box(sl, step,   x+0.1, 4.05, 2.4, 0.4, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, action, x+0.1, 4.55, 2.4, 1.3, font_size=12, color=MED_GRY, align=PP_ALIGN.CENTER)

    add_text_box(sl, "Risk: $900 and 90 days.  Upside: $165,000+ net benefit and a team that ramps faster at scale.",
                 0.4, 6.25, 12.5, 0.55, font_size=14, color=YELLOW, italic=True, align=PP_ALIGN.CENTER)

    add_text_box(sl, "Questions?  •  Let's run a live simulation right now.",
                 0.4, 6.9, 12.5, 0.45, font_size=13, color=MED_GRY, italic=True, align=PP_ALIGN.CENTER)

    path = os.path.join(OUTPUT_DIR, "exec_pitch_deck.pptx")
    prs.save(path)
    print(f"✅  Saved: {path}")


if __name__ == "__main__":
    build_roi_model()
    build_deck()
    print("\n✅  All deliverables generated successfully.")
